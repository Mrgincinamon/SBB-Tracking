"""Leitet aus dem kanonischen Foliendeck eine SOLO-Version ab.

Das kanonische Deck (presentation/SBB_Tracker_Praesentation.pptx, 22 Slides inkl.
manuell ergaenzter User-Story-Slides) ist fuer die Team-Aufnahme gedacht. Diese
Solo-Version ist optisch identisch, hat aber auf der Titel- und der Schluss-Slide
Sprechernotizen fuer eine Aufnahme durch eine Person (kein "wer spricht was").

Ausfuehren:  venv\\Scripts\\python.exe scripts\\make_solo_pptx.py
Output:      presentation/SBB_Tracker_Praesentation_solo.pptx
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).parent.parent
SRC = ROOT / "presentation" / "SBB_Tracker_Praesentation.pptx"
OUT = ROOT / "presentation" / "SBB_Tracker_Praesentation_solo.pptx"

SOLO_INTRO = (
    "SPRECHTEXT: «Grüezi, ich bin Joël Hasler. Ich stelle unser Gruppenprojekt zur "
    "Pünktlichkeit der SBB vor — eine datengetriebene Analyse von 2.74 Millionen "
    "echten Zug-Halten über 48 Tage. Ziel war, aus offenen Daten ehrliche, "
    "belastbare Aussagen abzuleiten.»\n\n"
    "HINWEIS (nur für dich): Du präsentierst durchgehend allein. In der "
    "Referentenansicht siehst du diese Notizen, das Publikum nicht. Achte auf die "
    "vorgegebene Gesamtlänge (Anzahl Studierende × 5 Min ≈ 10 Min).\n\n"
    "WORKFLOW: Vollbild präsentieren; beim Demo-Teil (Story-Schritte + Webapp) per "
    "Alt-Tab in die offene Live-App wechseln, klicken, dann zurück zu den Folien.\n\n"
    "TON: ruhig, nicht ablesen. Kernbotschaft: Effekte statistisch klar, aber "
    "praktisch klein — und du gehst ehrlich damit um."
)
SOLO_OUTRO = (
    "SPRECHTEXT: «Alle Daten sind offen, der gesamte Code liegt öffentlich auf "
    "GitHub und ist reproduzierbar. Vielen Dank fürs Zuschauen.»\n\n"
    "Ruhig abschliessen, kurze Pause am Ende lassen."
)


def main() -> int:
    if not SRC.exists():
        raise SystemExit(f"Kanonisches Deck fehlt: {SRC}")
    prs = Presentation(str(SRC))
    slides = list(prs.slides)
    slides[0].notes_slide.notes_text_frame.text = SOLO_INTRO
    slides[-1].notes_slide.notes_text_frame.text = SOLO_OUTRO
    prs.save(str(OUT))
    print(f"OK: Solo-Version erstellt {OUT}")
    print(f"     {len(prs.slides._sldIdLst)} Slides, {OUT.stat().st_size/1024:.0f} KB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
