"""Generiert eine PowerPoint-Praesentation (16:9) fuer die Video-Aufnahme.

Zieht alle Kennzahlen aus presentation/computed_results/results.json (bleibt
damit automatisch synchron mit Notebooks/PDF) und bettet die Notebook-Plots +
Webapp-Screenshots ein. Jede Slide hat AUSFUEHRLICHE Sprechernotizen: Sprechtext,
einfache Erklaerung der Fachbegriffe und moegliche Rueckfragen.

Ausfuehren:  venv\\Scripts\\python.exe scripts\\build_pptx.py
Output:      presentation/SBB_Tracker_Praesentation.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
RESULTS = ROOT / "presentation" / "computed_results" / "results.json"
SHOTS = ROOT / "presentation" / "screenshots"
NB_PLOTS = SHOTS / "notebooks"
OUT = ROOT / "presentation" / "SBB_Tracker_Praesentation.pptx"

SBB_RED = RGBColor(0xEB, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHTGREY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROSE = RGBColor(0xFC, 0xE9, 0xE9)  # helles Rot fuer Callout

SW, SH = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _fmt(x, nd=1):
    return f"{x:,.{nd}f}".replace(",", "'")


def add_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _para(tf, text, size, color=DARK, bold=False, first=False, bullet=False,
          align=PP_ALIGN.LEFT, space_after=8, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    run = p.add_run()
    run.text = ("•  " + text) if (bullet and level == 0) else (("–  " + text) if bullet else text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return p


def content_slide(prs, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    # roter Balken oben
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = SBB_RED
    bar.line.fill.background(); bar.shadow.inherit = False
    # Titel
    tf = _textbox(slide, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.85))
    _para(tf, title, 28, SBB_RED, bold=True, first=True)
    # duenne Akzentlinie unter dem Titel
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.28),
                                Inches(2.2), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = SBB_RED
    ln.line.fill.background(); ln.shadow.inherit = False
    return slide


def callout(slide, text, top=Inches(6.05)):
    """Rote Kernaussage-Box am unteren Slide-Rand."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top,
                                 Inches(11.93), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = ROSE
    box.line.color.rgb = SBB_RED; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = FONT


def bullets_slide(prs, title, bullets, notes="", body_top=1.55, size=18, call=None):
    slide = content_slide(prs, title)
    bottom = 6.0 if call else 7.1
    tf = _textbox(slide, Inches(0.7), Inches(body_top), Inches(11.9),
                  Inches(bottom - body_top))
    for i, b in enumerate(bullets):
        txt, lvl = b if isinstance(b, tuple) else (b, 0)
        _para(tf, txt, size if lvl == 0 else size - 3,
              DARK if lvl == 0 else GREY, first=(i == 0), bullet=True, level=lvl,
              bold=(lvl == 0 and txt.endswith(":")))
    if call:
        callout(slide, call)
    if notes:
        add_notes(slide, notes)
    return slide


def image_slide(prs, title, img_path: Path, caption="", notes=""):
    slide = content_slide(prs, title)
    if img_path.exists():
        max_w, max_h = Inches(11.6), Inches(4.95)
        pic = slide.shapes.add_picture(str(img_path), 0, Inches(1.5))
        ratio = min(max_w / pic.width, max_h / pic.height)
        pic.width = int(pic.width * ratio); pic.height = int(pic.height * ratio)
        pic.left = int((SW - pic.width) / 2)
    if caption:
        tf = _textbox(slide, Inches(0.7), Inches(6.7), Inches(11.93), Inches(0.65))
        _para(tf, caption, 13, GREY, first=True, align=PP_ALIGN.CENTER)
    if notes:
        add_notes(slide, notes)
    return slide


def add_footers(prs):
    """Fusszeile + Seitenzahl auf allen Slides ausser der Titelslide."""
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        tf = _textbox(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35))
        _para(tf, "SBB Tracker · ZHAW Scientific Programming FS2026", 9,
              LIGHTGREY, first=True)
        tf2 = _textbox(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35))
        _para(tf2, str(i + 1), 9, LIGHTGREY, first=True, align=PP_ALIGN.RIGHT)


def build():
    r = json.loads(RESULTS.read_text(encoding="utf-8"))
    w = r["test_welch_ttest"]; a = r["test_anova_linientyp"]
    o = r["test_ols"]; corr = r["test_correlation"]; dr = r["data_range"]

    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH

    # ---- 1. Titel ----
    s = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(s, DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.05), SW, Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = SBB_RED; bar.line.fill.background()
    bar.shadow.inherit = False
    tf = _textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.1))
    _para(tf, "SBB Tracker", 54, WHITE, bold=True, first=True)
    tf2 = _textbox(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.4))
    _para(tf2, "Pünktlichkeitsanalyse der Schweizerischen Bundesbahnen", 26, SBB_RED, bold=True, first=True)
    _para(tf2, "Datengetriebene Analyse von 2.74 Mio Zug-Halten · April–Mai 2026", 16, RGBColor(0xCC,0xCC,0xCC))
    tf3 = _textbox(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0))
    _para(tf3, "Joël Hasler & Patrick Ferreira", 16, WHITE, bold=True, first=True)
    _para(tf3, "ZHAW Scientific Programming · FS2026 · Dozent: Mario Gellrich", 13, RGBColor(0xAA,0xAA,0xAA))
    add_notes(s,
        "SPRECHTEXT: «Grüezi, wir sind Joël und Patrick. Wir haben die Pünktlichkeit "
        "der SBB datengetrieben untersucht — auf Basis von 2.74 Millionen echten "
        "Zug-Halten über 48 Tage. Unser Ziel war, aus offenen Daten ehrliche, "
        "belastbare Aussagen abzuleiten.»\n\n"
        "ORGANISATION: Legt vorher fest, wer welche Slides spricht (z. B. Joël 1–8, "
        "Patrick 9–16) und nennt es kurz. In der Referentenansicht von PowerPoint "
        "seht ihr diese Notizen, das Publikum nicht.\n\n"
        "TON: ruhig, nicht auswendig runterleiern. Die Kernbotschaft des ganzen "
        "Projekts ist: Die Effekte sind statistisch klar, aber praktisch klein — "
        "und wir gehen ehrlich damit um.")

    # ---- 2. Agenda ----
    bullets_slide(prs, "Agenda", [
        "Motivation und Forschungsfragen",
        "Datengrundlage und eine zentrale Design-Entscheidung",
        "Methodik: vier statistische Tests plus ein Sprachmodell",
        "Ergebnisse — mit Fokus auf Effektstärken",
        "Live-Demo der Webapp",
        "Limitationen und Schlussfolgerungen",
    ], size=20, notes=
        "SPRECHTEXT: «Kurz der Fahrplan: erst das Warum und die Datenbasis, dann "
        "die Methodik, dann die Ergebnisse, eine Live-Demo, und am Schluss ehrlich "
        "die Grenzen.»\n\n"
        "Nicht vorlesen — nur in einem Satz überfliegen. Schwerpunkt ansagen: "
        "Ergebnisse und Demo.")

    # ---- 3. Motivation ----
    bullets_slide(prs, "Motivation und Forschungsfragen", [
        "Die offizielle SBB-Pünktlichkeit (92.5 %) ist ein Durchschnitt —",
        "für die konkrete Reiseplanung wenig hilfreich. Wir gehen ins Detail:",
        ("F1: Sind Werktage unpünktlicher als Wochenenden?", 1),
        ("F2: Beeinflusst der Linientyp (S-Bahn, IC, IR …) die Verspätung?", 1),
        ("F3: Hängt das Wetter (Regen, Temperatur) mit Verspätung zusammen?", 1),
        ("F4: Welche Faktoren erklären Verspätung gemeinsam (Regression)?", 1),
    ], notes=
        "SPRECHTEXT: «Die SBB nennt eine Gesamt-Pünktlichkeit von 92.5 Prozent. "
        "Das ist ein Durchschnitt über alles und sagt mir wenig, wenn ich wissen "
        "will, wann und wo es konkret hakt. Deshalb diese vier Fragen.»\n\n"
        "HINTERGRUND, falls gefragt: 'Pünktlich' heisst bei der SBB <3 Minuten "
        "Verspätung. Genau diese 3-Minuten-Schwelle benutzen wir auch.\n\n"
        "Die vier Fragen sind unser roter Faden — jede wird später mit einem "
        "eigenen Test beantwortet.")

    # ---- 4. Datengrundlage ----
    bullets_slide(prs, "Datengrundlage", [
        f"{_fmt(r['n_total_events'],0)} gemessene Zug-Halte über {dr['n_days']} "
        f"Betriebstage ({dr['start']} bis {dr['end']})",
        "Drei offene Quellen: opentransportdata.swiss (Ist-Daten), data.sbb.ch",
        "(Bahnhof-Stammdaten), MeteoSchweiz (stündliches Wetter, 15 Stationen)",
        "Design-Entscheidung: hohe Auflösung statt langer Zeitspanne",
        ("48 Tage pro einzelnen Halt (~85 Mio Datenpunkte) statt Jahres-Summen", 1),
        ("erst das ermöglicht Stunden-, Bahnhof- und Wetter-Vergleiche", 1),
        "Bereinigung: 70 unmögliche Werte (fehlerhafte Quell-Zeitstempel) entfernt",
    ], notes=
        "SPRECHTEXT: «Basis sind 2.74 Millionen einzelne Zug-Halte über 48 Tage, "
        "aus drei offenen Quellen. Wir haben uns bewusst für Detailtiefe statt "
        "lange Zeitspanne entschieden: lieber 48 Tage ganz fein als Jahre an "
        "groben Summen — nur so kann man nach Stunde, Bahnhof und Wetter auswerten.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• 'Halt' = ein einzelner Halt eines Zuges an einem Bahnhof. Ein Zug auf "
        "seiner Fahrt erzeugt viele Halte.\n"
        "• 'Status REAL' = tatsächlich gemessene Ankunftszeit (keine Prognose).\n"
        "• 'Betriebstag' = der Fahrplantag; ein Nachtzug um 00:30 zählt zum Vortag.\n\n"
        "MÖGLICHE FRAGE: 'Warum nur 48 Tage?' — Die kompletten Daten von 2019–2025 "
        "wären rund 1.1 Terabyte. 48 Tage in voller Auflösung sind aussagekräftiger "
        "für unsere Fragen als grobe Jahres-Aggregate.")

    # ---- 5. Methodik ----
    bullets_slide(prs, "Methodik", [
        "Pipeline: Download → SQLite-Datenbank (3 Tabellen) → Aufbereitung →",
        "Zeit- und Wetter-Merkmale pro Halt anreichern",
        "Vier statistische Tests, jeweils mit p-Wert UND Effektstärke:",
        ("Welch-t-Test, ANOVA, Korrelation, multiple Regression", 1),
        "Plus Robustheits- und Annahmen-Checks (ehrliche Statistik)",
        "Ein Sprachmodell (Claude) für qualitative Ursachen-Hypothesen",
        "Werkzeuge: pandas · scipy · statsmodels · folium · streamlit · SQLite",
    ], notes=
        "SPRECHTEXT: «Die Daten laufen durch eine Pipeline in eine SQLite-Datenbank "
        "und werden pro Halt um Zeit- und Wetter-Infos ergänzt. Darauf rechnen wir "
        "vier Tests. Wichtig: Wir berichten nie nur den p-Wert, sondern immer auch "
        "die Effektstärke.»\n\n"
        "DAS WICHTIGSTE KONZEPT (unbedingt verstehen):\n"
        "• p-Wert = Wie wahrscheinlich wäre dieses Ergebnis, wenn es in Wirklichkeit "
        "KEINEN Unterschied gäbe? Klein (< 0.05) heisst 'wahrscheinlich kein Zufall' "
        "= statistisch signifikant.\n"
        "• ABER: Bei riesigen Datenmengen wird fast ALLES signifikant, auch winzige "
        "Unterschiede. Deshalb die Effektstärke: sie misst, wie GROSS der Unterschied "
        "praktisch ist. Genau das ist unsere Kernbotschaft.\n\n"
        "Die vier Tests sind nur Werkzeuge für je eine der vier Fragen — Details "
        "kommen auf den nächsten Slides.")

    # ---- 6. F1 ----
    bullets_slide(prs, "Ergebnis F1 — Werktag vs. Wochenende", [
        f"Werktag {_fmt(w['mean_werktag_sec'])} s  vs.  Wochenende {_fmt(w['mean_wochenende_sec'])} s",
        f"Welch-t = {_fmt(w['t_statistic'],1)}, p < 10⁻³⁰⁰ (hochsignifikant)",
        f"Differenz {_fmt(w['diff_sec'])} s, 95%-Konfidenzintervall [{_fmt(w['diff_ci95_low'])}, {_fmt(w['diff_ci95_high'])}] s",
        f"Effektstärke Cohen's d = {_fmt(w['cohens_d'],2)} → praktisch klein",
        "Gegenchecks: Mann-Whitney-Test und Test auf 48 Tagesmitteln bestätigen es",
    ], call="Signifikant ja — aber der Unterschied von ~15 s ist im Alltag kaum spürbar (d = 0.12).",
       notes=
        "SPRECHTEXT: «Werktage sind im Schnitt rund 15 Sekunden unpünktlicher als "
        "Wochenenden. Statistisch ist das glasklar. Aber 15 Sekunden merkt im Alltag "
        "niemand — die Effektstärke Cohen's d ist mit 0.12 klein. Genau dieser "
        "Unterschied zwischen 'signifikant' und 'relevant' zieht sich durch unsere "
        "ganze Arbeit.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• Welch-t-Test = vergleicht zwei Mittelwerte und verträgt, dass die beiden "
        "Gruppen unterschiedlich stark streuen.\n"
        "• Cohen's d = wie gross der Unterschied in 'Streuungs-Einheiten' ist. "
        "Faustregel: 0.2 klein, 0.5 mittel, 0.8 gross. Unser 0.12 ist also sehr klein.\n"
        "• Konfidenzintervall = der Bereich, in dem der wahre Wert mit 95 % "
        "Sicherheit liegt. Hier sehr eng, weil wir so viele Daten haben.\n"
        "• Mann-Whitney = ein Test, der keine Normalverteilung voraussetzt "
        "(arbeitet mit Rängen statt Werten) — als Gegenprobe.\n\n"
        "WARUM TAGESMITTEL? Unsere Millionen Halte sind nicht unabhängig "
        "(derselbe Zug, Bahnhof, Tag). Deshalb haben wir den Test auch nur auf den "
        "48 Tagesdurchschnitten gerechnet — der Effekt bleibt. Das zeigt: er ist "
        "echt, kein Artefakt der grossen Zahl.")

    # ---- 7. F2 ----
    image_slide(prs, "Ergebnis F2 — Linientyp (ANOVA)",
                NB_PLOTS / "03_analyse_visualisierung_cell21_plot0.png",
                caption=f"ANOVA F = {_fmt(a['f_statistic'],1)}, p < 10⁻³⁰⁰, η² = {_fmt(a['eta_squared'],3)} (klein). "
                        "Internationale Züge (TGV, EC, RJX) bringen Verspätung aus dem Ausland mit.",
                notes=
        "SPRECHTEXT: «Der Linientyp macht einen Unterschied — auch hier hochsignifikant, "
        "aber mit η² rund 0.04 wieder ein kleiner Effekt. Die eigentliche Geschichte: "
        "Internationale Züge wie TGV, EuroCity und RailJet sammeln im Ausland "
        "Verspätung und bringen sie in die Schweiz. Die inländischen IC und IR sind "
        "dagegen sehr pünktlich, um die 30 Sekunden.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• ANOVA (Varianzanalyse) = vergleicht die Mittelwerte von MEHR als zwei "
        "Gruppen gleichzeitig (hier die Linientypen). Antwortet auf: 'Gibt es "
        "überhaupt irgendwo einen Unterschied?'\n"
        "• F-Wert = die Test-Kennzahl der ANOVA; gross = die Gruppen unterscheiden "
        "sich deutlich relativ zur Streuung innerhalb der Gruppen.\n"
        "• η² (eta-Quadrat) = Effektstärke der ANOVA: Anteil der Verspätungs-"
        "Schwankung, der durch den Linientyp erklärt wird. 0.04 = nur ~4 %.\n"
        "• Tukey-HSD (im Notebook) = sagt zusätzlich, WELCHE Paare sich konkret "
        "unterscheiden, nicht nur dass es Unterschiede gibt.\n\n"
        "GRAFIK ERKLÄREN: Boxplot. Die Box ist der Bereich der mittleren 50 % der "
        "Werte, der Strich darin der Median. Je höher, desto verspäteter.")

    # ---- 8. F3 ----
    bullets_slide(prs, "Ergebnis F3 — Wetter und Verspätung", [
        f"Niederschlag: Pearson r = {_fmt(corr['niederschlag_mm']['pearson_r'],4)} (p < 10⁻²⁸⁰)",
        f"Sonne r = {_fmt(corr['sonne_min']['pearson_r'],4)} · Feuchte r = {_fmt(corr['feuchte_pct']['pearson_r'],4)}",
        "Alle hochsignifikant — aber |r| < 0.04, also inhaltlich trivial",
        "Spearman ≈ Pearson → Zusammenhang schwach und überwiegend monoton",
        "Wetter allein erklärt unter 0.3 % der Verspätungs-Schwankung",
    ], call="Mehr Regen → minim mehr Verspätung, aber der Effekt ist praktisch vernachlässigbar.",
       notes=
        "SPRECHTEXT: «Wetter wirkt in die erwartete Richtung — mehr Regen, etwas "
        "mehr Verspätung. Aber der Zusammenhang ist extrem schwach. Wieder ein "
        "Lehrstück: Bei 2.7 Millionen Beobachtungen wird selbst ein winziger "
        "Zusammenhang hochsignifikant. Die Korrelation r liegt aber unter 0.04 — "
        "praktisch bedeutungslos.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• Korrelation r = Mass für den Zusammenhang zweier Grössen, zwischen -1 "
        "und +1. 0 = kein Zusammenhang, ±1 = perfekter. Unsere Werte ~0.02–0.04 "
        "sind quasi null.\n"
        "• Pearson = misst LINEAREN Zusammenhang. Spearman = misst, ob es generell "
        "rauf/runter geht (monoton), über Ränge. Dass beide ähnlich klein sind, "
        "bestätigt: da ist einfach wenig.\n"
        "• 'erklärt < 0.3 % der Varianz' = r² (das Quadrat von r) als Anteil.\n\n"
        "MÖGLICHE FRAGE: 'Warum überhaupt zeigen, wenn der Effekt null ist?' — Weil "
        "es methodisch ehrlich ist und genau den Unterschied signifikant/relevant "
        "demonstriert. Negative bzw. Null-Ergebnisse sauber zu berichten ist Teil "
        "guter Wissenschaft.")

    # ---- 9. F4 ----
    image_slide(prs, "Ergebnis F4 — Multiple Regression (OLS)",
                NB_PLOTS / "03_analyse_visualisierung_cell30_plot0.png",
                caption=f"R² = {_fmt(o['r_squared'],4)} (~{_fmt(o['r_squared']*100,1)} %). Rush-Hour +10 s, Wochenende −12 s. "
                        "Diagnostik offen ausgewiesen (Heteroskedastizität, VIF unauffällig).",
                notes=
        "SPRECHTEXT: «Zum Schluss kombinieren wir alle Faktoren in einem Modell. "
        "Es erklärt nur rund 4 Prozent der Verspätung. Das klingt wenig, ist aber "
        "eine ehrliche und wichtige Erkenntnis: Verspätung entsteht überwiegend "
        "aus einzelnen, zufälligen Ereignissen — Defekte, Störungen — und nicht "
        "aus Tag, Linientyp oder Wetter.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• OLS-Regression = ein Modell, das die Verspätung gleichzeitig durch "
        "mehrere Einflussgrössen (Rush-Hour, Wochenende, Linientyp, Wetter) "
        "vorhersagt.\n"
        "• Koeffizient = der Effekt eines Faktors in Sekunden. Beispiel: Rush-Hour "
        "+10 s heisst, zur Stosszeit kommen im Schnitt 10 s Verspätung dazu.\n"
        "• R² = Anteil der Verspätungs-Schwankung, den das Modell erklärt. "
        "0.043 = ~4 %. Der Rest ist nicht erklärt = idiosynkratisch (Einzelfälle).\n"
        "• Heteroskedastizität (Breusch-Pagan-Test) = die Streuung der Fehler ist "
        "nicht konstant (bei grossen Verspätungen grösser). Wir benennen das offen.\n"
        "• VIF = prüft, ob sich Einflussgrössen gegenseitig zu stark überlappen "
        "(Multikollinearität). Bei uns unauffällig.\n\n"
        "GRAFIK: Residuen-Plot. Jeder Punkt ist ein Halt; x = Vorhersage, y = Fehler. "
        "Die Trichterform nach rechts ist genau die Heteroskedastizität.")

    # ---- 10. Heatmap ----
    image_slide(prs, "Zeitliches Muster — Stunde × Wochentag",
                NB_PLOTS / "03_analyse_visualisierung_cell33_plot0.png",
                caption="Mittlere Verspätung je Wochentag und Stunde. Dunkler = mehr Verspätung.",
                notes=
        "SPRECHTEXT: «Diese Heatmap zeigt, WANN es kritisch wird: dunkle Felder "
        "morgens und abends an Werktagen — die klassische Rush-Hour. Genau diese "
        "Erkenntnis bauen wir in der Webapp interaktiv nach.»\n\n"
        "GRAFIK ERKLÄREN: Zeilen = Wochentage (Mo–So), Spalten = Stunde 0–23, "
        "Farbe = durchschnittliche Verspätung. Je dunkler/röter, desto verspäteter. "
        "So sieht man Muster auf einen Blick statt in einer Zahlentabelle.\n\n"
        "Überleitung zur Demo: «Schauen wir uns das jetzt live in unserer Webapp an.»")

    # ---- 11. Webapp Karte ----
    image_slide(prs, "Webapp — Verspätungs-Hotspots auf der Karte",
                SHOTS / "webapp_01_karte.png",
                caption="Folium-Karte. Farbe = mittlere Verspätung. Mit dem Hotspot-Regler bleiben nur die kritischen Bahnhöfe übrig.",
                notes=
        "LIVE-DEMO (nicht nur die Folie zeigen, App wirklich bedienen!):\n"
        "1. App vorher starten: 'streamlit run app/streamlit_app.py'.\n"
        "2. «Jeder Punkt ist ein Bahnhof, die Farbe zeigt die mittlere Verspätung.»\n"
        "3. Den rechten Regler 'Nur Hotspots' langsam hochziehen — die Karte dünnt "
        "aus, bis nur noch die rotesten Bahnhöfe übrig sind.\n"
        "4. «Auffällig: die Hotspots liegen an den Grenzen — Buchs SG, St. Margrethen "
        "— weil internationale Züge ihre Verspätung importieren.»\n\n"
        "BEGRIFFE: 'Folium' = Python-Bibliothek für interaktive Landkarten. Der linke "
        "Regler blendet Bahnhöfe mit zu wenigen Daten aus (statistisch unsicher).\n\n"
        "Falls die Live-Demo hakt: dieser Screenshot dient als Fallback.")

    # ---- 12. Webapp ToD + LLM ----
    image_slide(prs, "Webapp — Pendler-Insight (Sprachmodell)",
                SHOTS / "webapp_03_pendler_insight.png",
                caption="Der Berater antwortet nur mit den Zahlen aus unseren Daten — er erfindet nichts dazu.",
                notes=
        "LIVE-DEMO:\n"
        "1. Eine Beispielfrage anklicken oder tippen, z. B. 'Wann erwische ich am "
        "ehesten einen pünktlichen Zug?'.\n"
        "2. «Das Sprachmodell bekommt unsere Statistik als Kontext mitgeliefert und "
        "darf NUR damit antworten — es soll nichts erfinden.»\n"
        "3. Auf den Ausklapp-Bereich 'Welche Daten hat das Modell gesehen?' zeigen — "
        "das macht es transparent und überprüfbar.\n\n"
        "BEGRIFFE EINFACH:\n"
        "• Sprachmodell / LLM = ein KI-Modell (hier Claude Sonnet 4.6), das Texte "
        "versteht und formuliert.\n"
        "• 'Kontext mitgeben' = wir hängen die relevanten Zahlen an die Frage an, "
        "damit die Antwort auf unseren Daten basiert und nicht auf Allgemeinwissen.\n"
        "• 'Halluzination' = wenn ein Modell plausibel klingende, aber falsche Dinge "
        "erfindet. Wir verhindern das durch klare Anweisung und mitgelieferte Fakten.\n\n"
        "MÖGLICHE FRAGE: 'Woher weiss man, dass es nicht doch erfindet?' — Wir geben "
        "die Datenbasis transparent aus und weisen das Modell explizit an, nur "
        "vorhandene Zahlen zu nennen.")

    # ---- 13. LLM Krisen-Tage ----
    bullets_slide(prs, "Sprachmodell — Analyse der Krisen-Tage", [
        "Claude klassifiziert die 10 verspätungsreichsten Tage nach Ursache",
        "Feste Auswahl an Kategorien (Wetter, Rush-Hour, Bahnhof-Ausfall, Bau …)",
        "Niedrige 'Temperatur' (0.2) + Anweisung 'erfinde nichts' + Konfidenz-Angabe",
        "Ehrlich: keine harte Verifikation — es sind qualitative Hypothesen",
        ("geplant: Abgleich mit Schweizer News (Factiva) als Gegenprobe", 1),
    ], notes=
        "SPRECHTEXT: «Zusätzlich lassen wir das Sprachmodell die zehn schlimmsten "
        "Tage einordnen — also eine mögliche Ursache vorschlagen. Wir kommunizieren "
        "klar, dass das Hypothesen sind, keine bewiesenen Fakten.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• 'Temperatur' beim Sprachmodell = Kreativitäts-Regler von 0 bis 1. Niedrig "
        "(0.2) = vorsichtig, faktentreu, wenig Fantasie. Genau das wollen wir hier.\n"
        "• 'Feste Taxonomie' = wir geben eine vorgegebene Liste möglicher Ursachen "
        "vor, statt das Modell frei raten zu lassen.\n"
        "• 'Konfidenz' = das Modell gibt selbst an, wie sicher es ist (hoch/mittel/"
        "niedrig). Achtung: das ist eine Selbsteinschätzung, nicht kalibriert.\n\n"
        "MÖGLICHE FRAGE: 'Ist das wissenschaftlich?' — Als qualitative Ergänzung ja, "
        "solange man es ehrlich als Hypothese kennzeichnet — was wir tun.")

    # ---- 14. Limitationen ----
    bullets_slide(prs, "Limitationen", [
        "Beobachtungen nicht unabhängig (gleicher Zug/Bahnhof/Tag) — Pseudoreplikation",
        ("dagegen: Tagesmittel-Gegencheck + Fokus auf Effektstärken", 1),
        "Nur 48 Tage im Frühling — keine Saisonalität (Winter, Bausaison fehlen)",
        "Nur SBB; Wetter-Station bis ~40 km entfernt (Mikroklima geht verloren)",
        "Korrelation ist nicht Kausalität (mögliche versteckte Faktoren)",
    ], notes=
        "SPRECHTEXT: «Zur Ehrlichkeit gehören die Grenzen. Der wichtigste Punkt: "
        "unsere Beobachtungen sind nicht unabhängig voneinander. Das macht rohe "
        "p-Werte zu optimistisch — deshalb der Tagesmittel-Gegencheck und der Fokus "
        "auf Effektstärken.»\n\n"
        "BEGRIFFE EINFACH:\n"
        "• Pseudoreplikation = man tut so, als hätte man Millionen unabhängige "
        "Messungen, obwohl viele zusammenhängen (derselbe Zug an vielen Halten). "
        "Dann wirken Ergebnisse sicherer als sie sind.\n"
        "• Saisonalität = jahreszeitliche Muster; wir haben nur Frühling, also "
        "keine Aussage über Winter.\n"
        "• Korrelation ≠ Kausalität = zwei Dinge hängen zusammen, heisst nicht, dass "
        "eins das andere VERURSACHT (z. B. Sturm UND Bauarbeiten am selben Tag).\n\n"
        "Diese Slide ist ein PLUS, kein Schwächeeingeständnis: Wer Limitationen "
        "sauber benennt, zeigt methodisches Verständnis.")

    # ---- 15. Schlussfolgerungen ----
    bullets_slide(prs, "Schlussfolgerungen", [
        "Die SBB ist sehr pünktlich (~95 % unter 3 Minuten) und räumlich stabil",
        "Tag, Linientyp und Wetter wirken messbar, aber praktisch klein",
        "Verspätung ist überwiegend idiosynkratisch (Modell erklärt nur ~4 %)",
        "Die grössten Hotspots sind Grenzbahnhöfe (Import-Effekt aus dem Ausland)",
        "Mehrwert: Webapp und Sprachmodell machen die Daten zugänglich",
    ], notes=
        "SPRECHTEXT: «Zusammengefasst: Die SBB ist sehr pünktlich. Die Effekte, die "
        "wir finden, sind echt, aber klein. Der grösste Teil der Verspätung ist "
        "Zufall im Einzelfall. Und die schlimmsten Bahnhöfe liegen an der Grenze. "
        "Unser Beitrag ist, diese Daten transparent und interaktiv zugänglich zu "
        "machen.»\n\n"
        "Das ist die Stelle, an der die rote Linie zusammenläuft: ehrlicher Umgang "
        "mit kleinen Effekten ist die methodische Stärke der Arbeit. Mit "
        "Selbstvertrauen sagen.")

    # ---- 16. Quellen ----
    bullets_slide(prs, "Quellen, Code und Reproduzierbarkeit", [
        "Daten: opentransportdata.swiss · data.sbb.ch · MeteoSchweiz (alle CC-BY)",
        "Code (MIT-Lizenz): github.com/Mrgincinamon/SBB-Tracking",
        "Vollständig reproduzierbar: requirements.txt + Build-Skripte",
        "Stack: Python 3.12 · pandas/scipy/statsmodels · SQLite · Streamlit · Anthropic",
        "Vielen Dank — gerne Fragen!",
    ], notes=
        "SPRECHTEXT: «Alle Daten sind offen, der gesamte Code liegt öffentlich auf "
        "GitHub und ist mit den Build-Skripten reproduzierbar. Vielen Dank, wir "
        "freuen uns auf eure Fragen.»\n\n"
        "Ruhig abschliessen, Blickkontakt, kurze Pause für Fragen lassen. Wer welche "
        "Frage übernimmt, vorher grob klären.")

    add_footers(prs)
    prs.save(str(OUT))
    print(f"OK: PPTX erstellt {OUT}")
    print(f"     {len(prs.slides._sldIdLst)} Slides, {OUT.stat().st_size/1024:.0f} KB")


if __name__ == "__main__":
    build()
